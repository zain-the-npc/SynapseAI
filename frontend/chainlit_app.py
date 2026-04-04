import os
import chainlit as cl
import requests
import base64
from pathlib import Path # Added for reliable path finding
from dotenv import load_dotenv # Added to actually load the file
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from supabase import create_client, Client

# 1. LOAD ENVIRONMENT VARIABLES 
# This finds the .env file in the SAME folder as this chainlit_app.py file
env_path = Path(__file__).parent / ".env"
load_dotenv(dotenv_path=env_path)

# 2. DATABASE & API CONFIGURATION
SUPABASE_URL = os.environ.get("SUPABASE_URL")
SUPABASE_KEY = os.environ.get("SUPABASE_KEY")
# Falls back to localhost if the environment variable isn't set yet
API_URL = os.environ.get("API_URL", "http://localhost:8000/chat")

if not SUPABASE_URL or not SUPABASE_KEY:
    print("CRITICAL ERROR: Supabase credentials missing in frontend/.env")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# 3. UPDATED AUTHENTICATION
@cl.password_auth_callback
async def auth(username: str, password: str):
    try:
        supabase.auth.sign_in_with_password({
            "email": username, 
            "password": password
        })
        return cl.User(identifier=username)
    except Exception as e:
        error_msg = str(e).lower()
        if "invalid login" in error_msg or "not found" in error_msg:
            try:
                signup_res = supabase.auth.sign_up({
                    "email": username, 
                    "password": password
                })
                # Check for error in response object
                if hasattr(signup_res, 'error') and signup_res.error:
                    raise Exception(signup_res.error.message)
                raise Exception("✅ Account Created! Please click Sign In again.")
            except Exception as signup_e:
                display_error = str(signup_e).replace("Exception:", "").strip()
                raise Exception(display_error)
        raise Exception("Invalid email or password. Please try again.")

# 4. SETTING UP STARTERS (Three peaceful buttons under search bar)
@cl.set_starters
async def set_starters():
    return [
        cl.Starter(
            label="📝 Make a Quiz",
            message="Create a 5-question quiz from the uploaded content.",
            icon="/public/quiz.svg",
        ),
        cl.Starter(
            label="📑 Summarize Paper",
            message="Provide a concise summary of the uploaded content.",
            icon="/public/summary.svg",
        ),
        cl.Starter(
            label="🗂️ Make Flashcards",
            message="Create study flashcards (Question/Answer) from this document.",
            icon="/public/cards.svg",
        ),
    ]

# 5. SILENT SESSION START
@cl.on_chat_start
async def start():
    user = cl.user_session.get("user")
    cl.user_session.set("active_image", None)
    cl.user_session.set("extracted_text", None)
    cl.user_session.set("file_type", None)
    cl.user_session.set("history", []) 
    
    try:
        response = supabase.table("chat_history") \
            .select("role, content") \
            .eq("user_id", user.identifier) \
            .order("created_at", desc=False) \
            .limit(10) \
            .execute()
        
        db_history = [{"role": row["role"], "content": row["content"]} for row in response.data]
        cl.user_session.set("history", db_history)
    except Exception as e:
        print(f"History Sync Error: {e}")

# 6. MESSAGE PROCESSING
@cl.on_message
async def main(message: cl.Message):
    history = cl.user_session.get("history")
    active_image = cl.user_session.get("active_image")
    user = cl.user_session.get("user")
    thread_id = cl.context.session.id 
    
    if message.elements:
        image_elements = [el for el in message.elements if "image" in el.mime]
        if image_elements:
            with open(image_elements[0].path, "rb") as f:
                active_image = base64.b64encode(f.read()).decode('utf-8')
                cl.user_session.set("active_image", active_image)
                cl.user_session.set("extracted_text", None)
                cl.user_session.set("file_type", "image")

        pdf_elements = [el for el in message.elements if "pdf" in el.mime or el.name.endswith('.pdf')]
        if pdf_elements:
            reader = PdfReader(pdf_elements[0].path)
            text = "".join([page.extract_text() + "\n" for page in reader.pages])
            cl.user_session.set("extracted_text", text)
            cl.user_session.set("active_image", None)
            cl.user_session.set("file_type", "pdf")

        pptx_elements = [el for el in message.elements if "presentation" in el.mime or el.name.endswith('.pptx')]
        if pptx_elements:
            prs = Presentation(pptx_elements[0].path)
            pptx_content = "".join([f"\n--- Slide {i+1} ---\n" + "".join([shape.text + "\n" for shape in slide.shapes if hasattr(shape, "text")]) for i, slide in enumerate(prs.slides)])
            cl.user_session.set("extracted_text", pptx_content)
            cl.user_session.set("active_image", None)
            cl.user_session.set("file_type", "pptx")

        docx_elements = [el for el in message.elements if "wordprocessingml" in el.mime or el.name.endswith('.docx')]
        if docx_elements:
            doc = Document(docx_elements[0].path)
            doc_content = "\n".join([para.text for para in doc.paragraphs])
            cl.user_session.set("extracted_text", doc_content)
            cl.user_session.set("active_image", None)
            cl.user_session.set("file_type", "docx")

    file_type = cl.user_session.get("file_type")
    loaders = {
        "image": "SynapseAI is analyzing your notes... 🧠",
        "pdf": "SynapseAI is processing the document... 📄",
        "pptx": "SynapseAI is extracting slides... 📊",
        "docx": "SynapseAI is reading the file... 📝",
        None: "SynapseAI is thinking... 💭"
    }
    
    msg = cl.Message(content="")
    loader_msg = loaders.get(file_type, "SynapseAI is thinking... 💭")
    await msg.stream_token(loader_msg)
    await msg.send()

    user_text = message.content if message.content else "Analyze this content."
    current_text = cl.user_session.get("extracted_text")
    final_message = user_text
    if current_text:
        label = file_type.upper() if file_type else "FILE"
        final_message = f"--- {label} CONTENT ---\n{current_text}\n--- END CONTENT ---\n\n{user_text}"

    payload = {
        "message": final_message,
        "history": history,
        "image_base64": active_image 
    }

    try:
        supabase.table("chat_history").insert({
            "user_id": user.identifier,
            "thread_id": thread_id,
            "role": "user",
            "content": user_text
        }).execute()
    except Exception as e:
        print(f"Error saving user message: {e}")

    try:
        full_answer = ""
        with requests.post(API_URL, json=payload, timeout=60, stream=True) as response:
            if response.status_code == 200:
                msg.content = "" 
                for chunk in response.iter_content(chunk_size=None, decode_unicode=True):
                    if chunk:
                        full_answer += chunk
                        await msg.stream_token(chunk)
                
                try:
                    supabase.table("chat_history").insert({
                        "user_id": user.identifier,
                        "thread_id": thread_id,
                        "role": "assistant",
                        "content": full_answer
                    }).execute()
                except Exception as e:
                    print(f"Error saving AI message: {e}")

                history.append({"role": "user", "content": user_text})
                history.append({"role": "assistant", "content": full_answer})
                cl.user_session.set("history", history[-10:])
                await msg.update()
            else:
                msg.content = f"❌ SynapseAI Error: Backend failure (Status: {response.status_code})."
                await msg.update()
    except Exception as e:
        msg.content = f"❌ SynapseAI Connection Error: {e}"
        await msg.update()